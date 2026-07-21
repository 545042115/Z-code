#!/usr/bin/env python3
"""Universal document processing CLI for the Ziner desktop app.

Supports extracting and creating DOCX, PPTX, XLSX, and PDF files.
This script is intended to be packaged into document-tools.exe with PyInstaller.
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Any


def read_json_arg(arg: str) -> dict[str, Any]:
    """Read JSON from a file path or a raw JSON string."""
    if arg.startswith('@'):
        with open(arg[1:], 'r', encoding='utf-8') as f:
            return json.load(f)
    if os.path.exists(arg):
        with open(arg, 'r', encoding='utf-8') as f:
            return json.load(f)
    return json.loads(arg)


def extract_docx(path: str) -> dict[str, Any]:
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        tables.append(rows)
    return {"type": "docx", "paragraphs": paragraphs, "tables": tables}


def extract_pdf(path: str) -> dict[str, Any]:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        return {"type": "pdf", "pages": pages, "text": "\n".join(pages)}
    except Exception as e:
        return {"type": "pdf", "error": str(e)}


def extract_pptx(path: str) -> dict[str, Any]:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append({"slide": i, "texts": texts})
    return {"type": "pptx", "slides": slides}


def extract_xlsx(path: str) -> dict[str, Any]:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    sheets = {}
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(cell) if cell is not None else "" for cell in row])
        sheets[name] = rows
    return {"type": "xlsx", "sheets": sheets}


def extract(path: str) -> dict[str, Any]:
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return extract_docx(path)
    if ext == ".pdf":
        return extract_pdf(path)
    if ext == ".pptx":
        return extract_pptx(path)
    if ext == ".xlsx":
        return extract_xlsx(path)
    return {"error": f"Unsupported extension: {ext}"}


def create_docx(output: str, spec: dict[str, Any]) -> None:
    from docx import Document
    from docx.shared import Inches, Pt

    doc = Document()
    title = spec.get("title")
    if title:
        doc.add_heading(title, level=0)
    for item in spec.get("content", []):
        typ = item.get("type", "paragraph")
        if typ == "heading":
            doc.add_heading(item.get("text", ""), level=item.get("level", 1))
        elif typ == "paragraph":
            doc.add_paragraph(item.get("text", ""))
        elif typ == "bullet":
            doc.add_paragraph(item.get("text", ""), style="List Bullet")
        elif typ == "numbered":
            doc.add_paragraph(item.get("text", ""), style="List Number")
        elif typ == "table":
            rows = item.get("rows", [])
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                for i, row in enumerate(rows):
                    for j, cell in enumerate(row):
                        table.cell(i, j).text = str(cell)
    doc.save(output)


def create_pptx(output: str, spec: dict[str, Any]) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title = spec.get("title")
    slides = spec.get("slides", [])
    if title and not slides:
        slides = [{"title": title, "content": []}]
    for s in slides:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = ""
        for line in s.get("content", []):
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
    prs.save(output)


def create_xlsx(output: str, spec: dict[str, Any]) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    default_sheet = wb.active
    sheets = spec.get("sheets", {"Sheet1": spec.get("data", [])})
    for idx, (name, rows) in enumerate(sheets.items()):
        if idx == 0:
            ws = default_sheet
            ws.title = name
        else:
            ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    wb.save(output)


def create_pdf(output: str, spec: dict[str, Any]) -> None:
    from fpdf import FPDF

    class PDF(FPDF):
        def header(self):
            pass

        def footer(self):
            self.set_y(-15)
            self.set_font("Arial", "I", 8)
            self.cell(0, 10, f"Page {self.page_no()}", align="C")

    pdf = PDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # FPDF default font Arial does not support CJK. For simple Latin text this is fine.
    # For Chinese content we would need a TrueType font with CJK support.
    pdf.set_font("Arial", "B", 16)
    title = spec.get("title", "")
    if title:
        pdf.cell(0, 10, title, ln=True)
        pdf.ln(5)

    pdf.set_font("Arial", "", 12)
    for item in spec.get("content", []):
        text = item.get("text", "") if isinstance(item, dict) else str(item)
        typ = item.get("type", "paragraph") if isinstance(item, dict) else "paragraph"
        if typ == "heading":
            pdf.set_font("Arial", "B", 14)
            pdf.cell(0, 10, text, ln=True)
            pdf.set_font("Arial", "", 12)
        else:
            pdf.multi_cell(0, 8, text)
        pdf.ln(2)
    pdf.output(output)


def create(tool: str, output: str, spec: dict[str, Any]) -> None:
    os.makedirs(os.path.dirname(os.path.abspath(output)) or ".", exist_ok=True)
    if tool == "docx":
        create_docx(output, spec)
    elif tool == "pptx":
        create_pptx(output, spec)
    elif tool == "xlsx":
        create_xlsx(output, spec)
    elif tool == "pdf":
        create_pdf(output, spec)
    else:
        raise ValueError(f"Unsupported create tool: {tool}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Universal document processing tool")
    subparsers = parser.add_subparsers(dest="command", required=True)

    extract_parser = subparsers.add_parser("extract", help="Extract text from a document")
    extract_parser.add_argument("input", help="Input file path")

    create_parser = subparsers.add_parser("create", help="Create a document from a JSON spec")
    create_parser.add_argument("tool", choices=["docx", "pptx", "xlsx", "pdf"], help="Document type")
    create_parser.add_argument("output", help="Output file path")
    create_parser.add_argument("--json", required=True, help="JSON spec (raw string or @path/to/spec.json)")

    args = parser.parse_args()

    try:
        if args.command == "extract":
            result = extract(args.input)
            print(json.dumps(result, ensure_ascii=False, indent=2))
        elif args.command == "create":
            spec = read_json_arg(args.json)
            create(args.tool, args.output, spec)
            print(json.dumps({"success": True, "output": args.output}, ensure_ascii=False))
        return 0
    except Exception as e:
        print(json.dumps({"error": str(e)}, ensure_ascii=False), file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
