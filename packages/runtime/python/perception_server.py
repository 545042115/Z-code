# @ziner/runtime — Python Perception Server
#
# Sidecar process providing OCR, image captioning, audio transcription,
# and document parsing. Uses ONNX Runtime for ML inference (no PyTorch).
# Communicates via stdin/stdout JSON-lines protocol.
#
# Protocol:
#   Request:  {"id": 1, "method": "ocr", "params": {"image": "<base64>"}}
#   Response: {"id": 1, "result": "extracted text"}
#   Error:    {"id": 1, "error": "error message"}
#
# Packaged with PyInstaller into perception-server.exe for distribution.

import sys
import json
import base64
import io
import os
import tempfile
import traceback
from pathlib import Path
from typing import Any


# ── Lazy-loaded backends ──────────────────────────────────────────────

_ocr_model = None
_asr_model = None


def get_ocr_model():
    """Initialize Surya OCR (ONNX-based, no PyTorch)."""
    global _ocr_model
    if _ocr_model is None:
        from surya.ocr import run_ocr
        from surya.model.recognition.model import load_model as load_rec_model
        from surya.model.recognition.processor import load_processor as load_rec_processor
        from surya.model.detection.model import load_model as load_det_model
        from surya.model.detection.processor import load_processor as load_det_processor

        det_model = load_det_model()
        det_processor = load_det_processor()
        rec_model = load_rec_model()
        rec_processor = load_rec_processor()

        _ocr_model = {
            'det_model': det_model,
            'det_processor': det_processor,
            'rec_model': rec_model,
            'rec_processor': rec_processor,
        }
    return _ocr_model


def get_asr_model():
    """Initialize faster-whisper (CTranslate2, no PyTorch)."""
    global _asr_model
    if _asr_model is None:
        from faster_whisper import WhisperModel
        model_size = os.environ.get('Z_ASR_MODEL', 'base')
        _asr_model = WhisperModel(model_size, device='cpu', compute_type='int8')
    return _asr_model


# ── Image utilities ───────────────────────────────────────────────────

def decode_image(image_b64: str) -> bytes:
    """Decode a base64 image to bytes."""
    return base64.b64decode(image_b64)


def get_image_metadata(image_bytes: bytes) -> dict:
    """Extract basic image metadata without ML."""
    from PIL import Image
    img = Image.open(io.BytesIO(image_bytes))
    info = {
        'width': img.width,
        'height': img.height,
        'format': img.format or 'unknown',
        'mode': img.mode,
    }
    # Get dominant colors (simple 5-color palette)
    if img.mode == 'RGBA':
        img = img.convert('RGBA')
    else:
        img = img.convert('RGB')
    small = img.resize((32, 32))
    colors = small.getcolors(1024)
    if colors:
        colors.sort(reverse=True)
        info['dominant_colors'] = [
            f'rgb({r},{g},{b})' for _, (r, g, b, *_) in colors[:5]
        ]
    return info


# ── Handlers ──────────────────────────────────────────────────────────

def handle_ocr(params: dict) -> str:
    """Extract text from a base64-encoded image using Surya OCR."""
    from PIL import Image
    image_b64 = params['image']
    langs = params.get('languages', ['en', 'zh'])
    image_bytes = decode_image(image_b64)
    img = Image.open(io.BytesIO(image_bytes)).convert('RGB')

    ocr = get_ocr_model()
    from surya.ocr import run_ocr
    predictions = run_ocr(
        [img],
        [langs],
        ocr['det_model'],
        ocr['det_processor'],
        ocr['rec_model'],
        ocr['rec_processor'],
    )

    lines = []
    for pred in predictions:
        for text_line in pred.text_lines:
            if text_line.text.strip():
                lines.append(text_line.text)
    return '\n'.join(lines)


def handle_caption(params: dict) -> str:
    """Describe an image using OCR + metadata (no heavy ML model needed)."""
    image_b64 = params['image']
    image_bytes = decode_image(image_b64)

    # Get metadata
    meta = get_image_metadata(image_bytes)

    # Try OCR for text content
    try:
        text = handle_ocr({'image': image_b64, 'languages': params.get('languages', ['en', 'zh'])})
    except Exception:
        text = ''

    parts = [
        f"This is a {meta['width']}x{meta['height']} {meta['format']} image.",
    ]
    if text.strip():
        parts.append(f"Text found in image:\n{text.strip()}")
    else:
        parts.append("No text detected in the image.")

    return '\n'.join(parts)


def handle_transcribe(params: dict) -> str:
    """Transcribe an audio file to text using faster-whisper."""
    audio_path = params['audio']
    model = get_asr_model()
    segments, _ = model.transcribe(audio_path, beam_size=5)
    text = ' '.join(seg.text for seg in segments)
    return text


def handle_parse_document(params: dict) -> str:
    """Parse a document file and return its text content.
    Supports: PDF, DOCX, PPTX, TXT
    """
    file_path = params['path']
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.pdf':
        import fitz
        doc = fitz.open(file_path)
        text = '\n'.join(page.get_text() for page in doc)
        doc.close()
        return text

    elif ext == '.docx':
        import docx
        doc = docx.Document(file_path)
        return '\n'.join(p.text for p in doc.paragraphs)

    elif ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text)
        return '\n'.join(texts)

    elif ext == '.txt':
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()

    else:
        raise ValueError(f'Unsupported document type: {ext}')


def handle_ocr_file(params: dict) -> str:
    """OCR a PDF file page by page."""
    from PIL import Image
    file_path = params['path']
    langs = params.get('languages', ['en', 'zh'])
    import fitz
    doc = fitz.open(file_path)
    texts = []

    ocr = get_ocr_model()
    from surya.ocr import run_ocr

    for page_num, page in enumerate(doc):
        pix = page.get_pixmap(dpi=200)
        img_bytes = pix.tobytes('png')
        img = Image.open(io.BytesIO(img_bytes)).convert('RGB')
        predictions = run_ocr(
            [img], [langs],
            ocr['det_model'], ocr['det_processor'],
            ocr['rec_model'], ocr['rec_processor'],
        )
        page_lines = []
        for pred in predictions:
            for text_line in pred.text_lines:
                if text_line.text.strip():
                    page_lines.append(text_line.text)
        if page_lines:
            texts.append(f'--- Page {page_num + 1} ---\n' + '\n'.join(page_lines))

    doc.close()
    return '\n\n'.join(texts)


# ── Dispatch ──────────────────────────────────────────────────────────

HANDLERS: dict[str, Any] = {
    'ocr': handle_ocr,
    'caption': handle_caption,
    'transcribe': handle_transcribe,
    'parse_document': handle_parse_document,
    'ocr_file': handle_ocr_file,
}


def handle_request(req: dict) -> dict:
    method = req.get('method')
    params = req.get('params', {})
    req_id = req.get('id', 0)

    if method not in HANDLERS:
        return {'id': req_id, 'error': f'Unknown method: {method}'}

    try:
        result = HANDLERS[method](params)
        return {'id': req_id, 'result': result}
    except Exception as e:
        traceback.print_exc()
        return {'id': req_id, 'error': str(e)}


def main():
    """Read JSON lines from stdin, write responses to stdout."""
    sys.stderr.write('[perception_server] ready\n')
    sys.stderr.flush()

    for line in sys.stdin:
        line = line.strip()
        if not line:
            continue
        try:
            req = json.loads(line)
            resp = handle_request(req)
        except json.JSONDecodeError as e:
            resp = {'id': 0, 'error': f'Invalid JSON: {e}'}

        sys.stdout.write(json.dumps(resp, ensure_ascii=False) + '\n')
        sys.stdout.flush()


if __name__ == '__main__':
    main()
